"""
Notion to PowerPoint (Simple Mode)
==================================

Directly converts any Notion page to PowerPoint slides.
No LLM analysis - just preserves the original structure.

Usage:
    from teams.notion_to_pptx import notion_to_pptx

    result = notion_to_pptx(
        notion_url="https://www.notion.so/your-page-id",
        title="My Presentation"  # optional, uses page title if not provided
    )
    print(result["powerpoint_path"])
"""

import os
import re
import logging
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime

from dotenv import load_dotenv
load_dotenv()

from pptx import Presentation
from pptx.util import Inches, Pt

from agents.notion_reader import read_notion_page

# ============================================================================
# Logging Setup
# ============================================================================
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("NotionToPPTX")

# ============================================================================
# Constants
# ============================================================================
OUTPUT_DIR = Path(os.getenv("OUTPUT_DIR", "./output"))
OUTPUT_DIR.mkdir(exist_ok=True)

TEMPLATE_PATH = Path(os.getenv(
    "PPTX_TEMPLATE_PATH",
    str(Path(__file__).parent.parent / "templates" / "sweetspot_template.pptx")
))


# ============================================================================
# Content Parser
# ============================================================================
def parse_notion_content(content: str) -> List[Dict[str, Any]]:
    """
    Parse Notion markdown content into sections.

    Returns list of sections, each with:
    - title: section heading
    - level: heading level (1, 2, 3)
    - bullets: list of bullet points
    - paragraphs: list of paragraph text
    """
    sections = []
    current_section = None

    lines = content.split('\n')

    for line in lines:
        line_stripped = line.strip()

        if not line_stripped:
            continue

        # Check for headings
        if line_stripped.startswith('# '):
            # Save previous section
            if current_section:
                sections.append(current_section)
            current_section = {
                'title': line_stripped[2:].strip(),
                'level': 1,
                'bullets': [],
                'paragraphs': []
            }
        elif line_stripped.startswith('## '):
            if current_section:
                sections.append(current_section)
            current_section = {
                'title': line_stripped[3:].strip(),
                'level': 2,
                'bullets': [],
                'paragraphs': []
            }
        elif line_stripped.startswith('### '):
            if current_section:
                sections.append(current_section)
            current_section = {
                'title': line_stripped[4:].strip(),
                'level': 3,
                'bullets': [],
                'paragraphs': []
            }
        elif line_stripped.startswith('- ') or line_stripped.startswith('* '):
            # Bullet point
            if current_section is None:
                current_section = {
                    'title': 'Content',
                    'level': 1,
                    'bullets': [],
                    'paragraphs': []
                }
            current_section['bullets'].append(line_stripped[2:].strip())
        elif re.match(r'^\d+\.\s', line_stripped):
            # Numbered list
            if current_section is None:
                current_section = {
                    'title': 'Content',
                    'level': 1,
                    'bullets': [],
                    'paragraphs': []
                }
            # Remove the number prefix
            text = re.sub(r'^\d+\.\s*', '', line_stripped)
            current_section['bullets'].append(text)
        else:
            # Regular paragraph
            if current_section is None:
                current_section = {
                    'title': 'Content',
                    'level': 1,
                    'bullets': [],
                    'paragraphs': []
                }
            # Skip very short lines (likely formatting artifacts)
            if len(line_stripped) > 3:
                current_section['paragraphs'].append(line_stripped)

    # Don't forget the last section
    if current_section:
        sections.append(current_section)

    return sections


# ============================================================================
# PowerPoint Generation
# ============================================================================
def create_pptx(
    title: str,
    sections: List[Dict[str, Any]],
    output_name: str
) -> Tuple[Optional[Path], int]:
    """
    Create PowerPoint from parsed sections.

    Returns (path, slide_count) or (None, 0) on failure.
    """
    try:
        # Load template or create blank
        if TEMPLATE_PATH.exists():
            prs = Presentation(str(TEMPLATE_PATH))
            # Clear existing slides
            while len(prs.slides) > 0:
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]
        else:
            prs = Presentation()

        slides_created = 0

        # ================================================================
        # Title Slide
        # ================================================================
        layout_idx = min(0, len(prs.slide_layouts) - 1)
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        slides_created += 1

        # Add title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True

        # Add date
        p2 = tf.add_paragraph()
        p2.text = datetime.now().strftime("%B %Y")
        p2.font.size = Pt(20)

        # ================================================================
        # Agenda Slide (if multiple sections)
        # ================================================================
        if len(sections) > 1:
            layout_idx = min(1, len(prs.slide_layouts) - 1)
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            slides_created += 1

            if slide.shapes.title:
                slide.shapes.title.text = "Agenda"

            # Find content placeholder or create textbox
            content_shape = None
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    content_shape = shape
                    break

            if content_shape:
                tf = content_shape.text_frame
                tf.clear()
                for i, section in enumerate(sections):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = f"{i+1}. {section['title']}"
                    p.font.size = Pt(18)
            else:
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
                tf = txBox.text_frame
                for i, section in enumerate(sections):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = f"{i+1}. {section['title']}"
                    p.font.size = Pt(18)

        # ================================================================
        # Content Slides
        # ================================================================
        for section in sections:
            # Combine bullets and paragraphs
            all_content = section['bullets'] + section['paragraphs']

            if not all_content:
                # Empty section - just create a divider slide
                layout_idx = min(2, len(prs.slide_layouts) - 1)
                slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
                slides_created += 1

                if slide.shapes.title:
                    slide.shapes.title.text = section['title']
                else:
                    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(2))
                    tf = txBox.text_frame
                    p = tf.paragraphs[0]
                    p.text = section['title']
                    p.font.size = Pt(36)
                    p.font.bold = True
                continue

            # Split content into chunks of 6 items per slide
            chunks = [all_content[i:i+6] for i in range(0, len(all_content), 6)]

            for chunk_idx, chunk in enumerate(chunks):
                layout_idx = min(1, len(prs.slide_layouts) - 1)
                slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
                slides_created += 1

                # Title
                slide_title = section['title']
                if chunk_idx > 0:
                    slide_title = f"{section['title']} (cont.)"

                if slide.shapes.title:
                    slide.shapes.title.text = slide_title

                # Content
                content_shape = None
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        content_shape = shape
                        break

                if content_shape:
                    tf = content_shape.text_frame
                    tf.clear()
                    for i, text in enumerate(chunk):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = f"• {text}"
                        p.font.size = Pt(16)
                else:
                    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
                    tf = txBox.text_frame
                    for i, text in enumerate(chunk):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = f"• {text}"
                        p.font.size = Pt(16)

        # Save
        safe_name = "".join(c if c.isalnum() else "_" for c in output_name)
        pptx_filename = f"notion_pptx_{safe_name}.pptx"
        pptx_path = OUTPUT_DIR / pptx_filename
        prs.save(str(pptx_path))

        return pptx_path, slides_created

    except Exception as e:
        logger.error(f"PowerPoint generation failed: {e}")
        return None, 0


# ============================================================================
# Main Entry Point
# ============================================================================
def notion_to_pptx(
    notion_url: str = None,
    raw_content: str = None,
    title: str = None,
    output_name: str = None
) -> Dict[str, Any]:
    """
    Convert Notion page directly to PowerPoint.

    No LLM analysis - preserves original structure.

    Args:
        notion_url: Notion page URL
        raw_content: Raw markdown content (alternative to notion_url)
        title: Presentation title (optional, uses page title if not provided)
        output_name: Output filename base (optional)

    Returns:
        Dict with:
        - success: bool
        - powerpoint_path: path to generated .pptx
        - slide_count: number of slides
        - sections: list of section titles
        - error: error message if failed
    """
    result = {
        "success": False,
        "powerpoint_path": None,
        "slide_count": 0,
        "sections": [],
        "error": None
    }

    logger.info("=" * 60)
    logger.info("NOTION TO PPTX (Simple Mode)")
    logger.info("=" * 60)

    # ========================================================================
    # Step 1: Get content
    # ========================================================================
    page_title = title or "Presentation"

    if notion_url:
        logger.info(f"Reading Notion page: {notion_url}")
        notion_result = read_notion_page(notion_url)

        if not notion_result["success"]:
            result["error"] = f"Failed to read Notion: {notion_result['error']}"
            logger.error(result["error"])
            return result

        content = notion_result["content"]
        page_title = title or notion_result.get("title", "Presentation")
        logger.info(f"Page title: {page_title}")

    elif raw_content:
        content = raw_content

    else:
        result["error"] = "No content provided (need notion_url or raw_content)"
        return result

    logger.info(f"Content length: {len(content)} characters")

    # ========================================================================
    # Step 2: Parse content into sections
    # ========================================================================
    logger.info("Parsing content structure...")
    sections = parse_notion_content(content)

    logger.info(f"Found {len(sections)} sections:")
    for s in sections:
        logger.info(f"  - {s['title']}: {len(s['bullets'])} bullets, {len(s['paragraphs'])} paragraphs")

    result["sections"] = [s['title'] for s in sections]

    # ========================================================================
    # Step 3: Generate PowerPoint
    # ========================================================================
    logger.info("Generating PowerPoint...")

    output_name = output_name or page_title
    pptx_path, slide_count = create_pptx(page_title, sections, output_name)

    if pptx_path:
        result["success"] = True
        result["powerpoint_path"] = str(pptx_path)
        result["slide_count"] = slide_count
        logger.info(f"Created: {pptx_path} ({slide_count} slides)")
    else:
        result["error"] = "PowerPoint generation failed"

    logger.info("=" * 60)
    logger.info("DONE")
    logger.info("=" * 60)

    return result


# ============================================================================
# Main
# ============================================================================
if __name__ == "__main__":
    # Test with sample content
    test_content = """
# My Project Plan

## Overview
This is a project about building something cool.
We want to deliver value to customers.

## Goals
- Launch MVP by Q2
- Get 1000 users in first month
- Achieve 4.5 star rating

## Timeline
1. Design phase - January
2. Development - February/March
3. Testing - April
4. Launch - May

## Team
- Alice - Product Manager
- Bob - Lead Developer
- Carol - Designer

## Risks
- Timeline might slip
- Budget constraints
- Technical complexity
"""

    print("Testing Notion to PPTX (Simple Mode)...")

    result = notion_to_pptx(
        raw_content=test_content,
        title="My Project Plan"
    )

    print(f"\nSuccess: {result['success']}")
    print(f"Slides: {result['slide_count']}")
    print(f"Sections: {result['sections']}")
    print(f"PowerPoint: {result['powerpoint_path']}")
