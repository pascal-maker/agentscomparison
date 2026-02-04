"""
Smart Discovery Workflow (Production Grade)
===========================================

Evidence-grounded discovery report generation from any Notion page.

Pipeline:
1. NotionReader → EvidenceItem[] (with IDs, quotes, paths)
2. ContentAnalyzer → sections + evidence mapping
3. CustomerConfig → enforce must_include, terminology, limits
4. Revisor → enforce grounding, add Open Questions for gaps
5. PowerPointWriter → Sweetspot deck (Title, Agenda, Sections)

Every bullet in output MUST have [EVID-xxx] reference.
"""

import os
import re
import logging
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime

from dotenv import load_dotenv
load_dotenv()

from pptx import Presentation
from pptx.util import Inches, Pt

from agents.notion_reader import read_notion_page
from agents.content_analyzer import analyze_with_evidence
from agents.revisor import revise_sections, enforce_slide_budget
from shared.evidence import (
    EvidenceItem,
    EvidenceCollection,
    GroundedSection,
    GroundedBullet,
    CustomerConfig,
    extract_evidence_from_content,
    validate_grounded_report
)

# ============================================================================
# Logging Setup
# ============================================================================
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("SmartDiscovery")

# ============================================================================
# Constants
# ============================================================================
OUTPUT_DIR = Path(os.getenv("OUTPUT_DIR", "./output"))
OUTPUT_DIR.mkdir(exist_ok=True)

TEMPLATE_PATH = Path(os.getenv(
    "PPTX_TEMPLATE_PATH",
    str(Path(__file__).parent.parent / "templates" / "sweetspot_template.pptx")
))

# Slide layout indices in Sweetspot template
LAYOUT_BLANK = 0      # Title/cover slides
LAYOUT_CONTENT = 1    # Bullet point content
LAYOUT_CHAPTER = 2    # Section dividers
LAYOUT_TEXT_LEFT = 3  # Text on left
LAYOUT_TEXT_RIGHT = 4 # Text on right


# ============================================================================
# Main Entry Point
# ============================================================================
def smart_discover(
    notion_url: str = None,
    raw_content: str = None,
    customer_name: str = "Client",
    config: CustomerConfig = None,
) -> Dict[str, Any]:
    """
    Production-grade smart discovery with evidence grounding.

    Args:
        notion_url: Notion page URL to read (optional)
        raw_content: Raw content string (alternative to notion_url)
        customer_name: Customer name for the report
        config: Optional CustomerConfig for constraints

    Returns:
        Dict with:
        - success: bool
        - markdown_path: path to generated .md file
        - powerpoint_path: path to generated .pptx file
        - sections: list of section names
        - evidence_count: number of evidence items
        - validation: grounding validation results
        - error: error message if failed
    """
    result = {
        "success": False,
        "markdown_path": None,
        "powerpoint_path": None,
        "sections": [],
        "evidence_count": 0,
        "validation": None,
        "error": None,
    }

    # Create default config if not provided
    if not config:
        config = CustomerConfig(name=customer_name)

    logger.info("=" * 60)
    logger.info("SMART DISCOVERY WORKFLOW (Production)")
    logger.info("=" * 60)
    logger.info(f"Customer: {config.name}")

    # ========================================================================
    # Step 1: Get content and extract evidence
    # ========================================================================
    logger.info("\n[Step 1] Extracting evidence from input...")

    page_title = ""
    page_id = ""
    page_url = notion_url or ""

    if notion_url:
        logger.info(f"Reading Notion page: {notion_url}")
        notion_result = read_notion_page(notion_url)
        if not notion_result["success"]:
            result["error"] = f"Failed to read Notion: {notion_result['error']}"
            return result
        content = notion_result["content"]
        page_title = notion_result.get("title", "")
        page_id = notion_result.get("metadata", {}).get("page_id", "")
    elif raw_content:
        content = raw_content
        page_title = "Input Document"
    else:
        result["error"] = "No content provided (need notion_url or raw_content)"
        return result

    # Extract evidence items
    evidence = extract_evidence_from_content(
        content=content,
        page_title=page_title,
        page_id=page_id,
        page_url=page_url
    )

    result["evidence_count"] = len(evidence)
    logger.info(f"Extracted {len(evidence)} evidence items")

    # ========================================================================
    # Step 2: Analyze content and discover sections
    # ========================================================================
    logger.info("\n[Step 2] Analyzing content to discover sections...")

    sections, title, summary = analyze_with_evidence(evidence, config)

    logger.info(f"Discovered title: {title}")
    logger.info(f"Discovered {len(sections)} sections")

    # ========================================================================
    # Step 3: Revisor pass - enforce grounding rules
    # ========================================================================
    logger.info("\n[Step 3] Revisor pass - enforcing grounding rules...")

    revised_sections, revision_result = revise_sections(sections, evidence, config)

    logger.info(revision_result.summary())

    # ========================================================================
    # Step 4: Enforce slide budget
    # ========================================================================
    logger.info("\n[Step 4] Enforcing slide budget...")

    final_sections = enforce_slide_budget(revised_sections, config)

    # ========================================================================
    # Step 5: Validate grounding
    # ========================================================================
    logger.info("\n[Step 5] Validating evidence grounding...")

    validation = validate_grounded_report(final_sections)
    result["validation"] = validation

    if not validation["valid"]:
        logger.warning(f"Grounding validation failed: {validation['errors']}")
        # Continue anyway - we'll have Open Questions for gaps

    result["sections"] = [s.name for s in final_sections]

    # ========================================================================
    # Step 6: Generate Markdown output
    # ========================================================================
    logger.info("\n[Step 6] Generating Markdown report...")

    markdown_content = _generate_markdown(
        title=title,
        summary=summary,
        sections=final_sections,
        evidence=evidence,
        config=config
    )

    # Save markdown
    safe_name = "".join(c if c.isalnum() else "_" for c in config.name)
    md_filename = f"smart_discovery_{safe_name}.md"
    md_path = OUTPUT_DIR / md_filename
    md_path.write_text(markdown_content)
    result["markdown_path"] = str(md_path)
    logger.info(f"Saved markdown: {md_path}")

    # ========================================================================
    # Step 7: Generate PowerPoint output
    # ========================================================================
    logger.info("\n[Step 7] Generating PowerPoint presentation...")

    pptx_path, slide_count = _generate_powerpoint(
        title=title,
        summary=summary,
        sections=final_sections,
        evidence=evidence,
        config=config
    )

    if pptx_path:
        result["powerpoint_path"] = str(pptx_path)
        logger.info(f"Saved PowerPoint: {pptx_path} ({slide_count} slides)")
    else:
        logger.warning("PowerPoint generation failed")

    result["success"] = True

    logger.info("\n" + "=" * 60)
    logger.info("SMART DISCOVERY COMPLETE")
    logger.info("=" * 60)

    return result


# ============================================================================
# Markdown Generation
# ============================================================================
def _generate_markdown(
    title: str,
    summary: str,
    sections: List[GroundedSection],
    evidence: EvidenceCollection,
    config: CustomerConfig,
) -> str:
    """Generate markdown report with evidence references."""

    lines = [
        f"# {title}",
        "",
        f"**Customer:** {config.name}",
        f"**Generated:** {datetime.now().strftime('%Y-%m-%d %H:%M')}",
        f"**Evidence Items:** {len(evidence)}",
        "",
        "---",
        "",
    ]

    # Executive Summary
    lines.extend([
        "## Executive Summary",
        "",
        summary if summary else "*No summary available - see sections below.*",
        "",
        "---",
        "",
    ])

    # Agenda / Table of Contents
    lines.extend([
        "## Agenda",
        "",
    ])
    for i, section in enumerate(sections, 1):
        lines.append(f"{i}. {section.name}")
    lines.extend(["", "---", ""])

    # Sections with evidence references
    for section in sections:
        lines.append(f"## {section.name}")
        lines.append("")

        if section.description:
            lines.append(section.description)
            lines.append("")

        # Grounded bullets with evidence IDs
        if section.bullets:
            for bullet in section.bullets:
                refs = " ".join(f"[{eid}]" for eid in bullet.evidence_ids)
                lines.append(f"- {bullet.text} {refs}")
            lines.append("")

        # Open Questions
        if section.open_questions:
            lines.append("**Open Questions:**")
            for q in section.open_questions:
                lines.append(f"- {q}")
            lines.append("")

        # Evidence footer for this section
        if section.evidence_ids:
            lines.append("**Evidence Sources:**")
            for eid in set(section.evidence_ids):
                item = evidence.get(eid)
                if item:
                    lines.append(f"- {item.format_citation()}")
            lines.append("")

        lines.append("---")
        lines.append("")

    # Full evidence appendix
    lines.extend([
        "## Evidence Appendix",
        "",
        "Complete list of evidence items referenced in this report:",
        "",
    ])

    for item in evidence.items.values():
        lines.append(f"- **{item.id}**: {item.quote[:100]}{'...' if len(item.quote) > 100 else ''}")

    lines.extend([
        "",
        "---",
        "",
        f"*Report generated by Smart Discovery on {datetime.now().isoformat()}*",
    ])

    return "\n".join(lines)


# ============================================================================
# PowerPoint Generation
# ============================================================================
def _generate_powerpoint(
    title: str,
    summary: str,
    sections: List[GroundedSection],
    evidence: EvidenceCollection,
    config: CustomerConfig,
) -> Tuple[Optional[Path], int]:
    """
    Generate Sweetspot-styled PowerPoint with evidence in speaker notes.

    Returns:
        Tuple of (output_path, slide_count) or (None, 0) on failure
    """
    try:
        # Load template
        if TEMPLATE_PATH.exists():
            prs = Presentation(str(TEMPLATE_PATH))
            # Clear existing slides
            while len(prs.slides) > 0:
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]
        else:
            logger.warning(f"Template not found: {TEMPLATE_PATH}, using blank")
            prs = Presentation()

        slides_created = 0
        max_slides = config.get_max_slides()

        # ====================================================================
        # Slide 1: Title Slide
        # ====================================================================
        layout_idx = min(LAYOUT_BLANK, len(prs.slide_layouts) - 1)
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        slides_created += 1

        # Add title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True

        # Add customer name
        p2 = tf.add_paragraph()
        p2.text = config.name
        p2.font.size = Pt(28)

        # Add date
        p3 = tf.add_paragraph()
        p3.text = datetime.now().strftime("%B %Y")
        p3.font.size = Pt(18)

        # Speaker notes with metadata
        notes = slide.notes_slide
        notes.notes_text_frame.text = f"Title slide\nEvidence items: {len(evidence)}\nGenerated: {datetime.now().isoformat()}"

        # ====================================================================
        # Slide 2: Agenda Slide
        # ====================================================================
        layout_idx = min(LAYOUT_CONTENT, len(prs.slide_layouts) - 1)
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        slides_created += 1

        if slide.shapes.title:
            slide.shapes.title.text = "Agenda"

        # Add section list
        content_shape = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                content_shape = shape
                break

        if content_shape:
            tf = content_shape.text_frame
            tf.clear()
            for i, section in enumerate(sections):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"{i+1}. {section.name}"
                p.font.size = Pt(18)
        else:
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
            tf = txBox.text_frame
            for i, section in enumerate(sections):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"{i+1}. {section.name}"
                p.font.size = Pt(18)

        notes = slide.notes_slide
        notes.notes_text_frame.text = f"Agenda slide\nSections: {len(sections)}"

        # ====================================================================
        # Section Slides
        # ====================================================================
        per_section_max = config.get_per_section_max()

        for section in sections:
            if slides_created >= max_slides:
                logger.warning(f"Reached slide budget ({max_slides}), stopping")
                break

            # Section divider slide
            layout_idx = min(LAYOUT_CHAPTER, len(prs.slide_layouts) - 1)
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            slides_created += 1

            if slide.shapes.title:
                slide.shapes.title.text = section.name
            else:
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(2))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = section.name
                p.font.size = Pt(36)
                p.font.bold = True

            notes = slide.notes_slide
            notes.notes_text_frame.text = f"Section: {section.name}\nBullets: {len(section.bullets)}\nOpen Questions: {len(section.open_questions)}"

            # Content slides (max 6 bullets per slide)
            all_items = []
            for bullet in section.bullets:
                refs = " ".join(f"[{eid}]" for eid in bullet.evidence_ids)
                all_items.append((f"{bullet.text} {refs}", bullet.evidence_ids))

            for q in section.open_questions:
                all_items.append((f"[OPEN] {q}", []))

            # Split into chunks of 6
            chunks = [all_items[i:i+6] for i in range(0, len(all_items), 6)]

            section_slides = 0
            for chunk_idx, chunk in enumerate(chunks):
                if slides_created >= max_slides:
                    break
                if section_slides >= per_section_max:
                    break

                layout_idx = min(LAYOUT_CONTENT, len(prs.slide_layouts) - 1)
                slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
                slides_created += 1
                section_slides += 1

                # Title
                slide_title = section.name if chunk_idx == 0 else f"{section.name} (cont.)"
                if slide.shapes.title:
                    slide.shapes.title.text = slide_title

                # Content
                content_shape = None
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        content_shape = shape
                        break

                if content_shape:
                    tf = content_shape.text_frame
                    tf.clear()
                    for i, (text, _) in enumerate(chunk):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = text
                        p.font.size = Pt(16)
                else:
                    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
                    tf = txBox.text_frame
                    for i, (text, _) in enumerate(chunk):
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = "• " + text
                        p.font.size = Pt(16)

                # Speaker notes with full evidence details
                notes_lines = [f"Section: {section.name}", "Evidence details:"]
                for text, eids in chunk:
                    for eid in eids:
                        item = evidence.get(eid)
                        if item:
                            notes_lines.append(f"  {eid}: {item.quote[:80]}...")
                notes = slide.notes_slide
                notes.notes_text_frame.text = "\n".join(notes_lines)

        # Save
        safe_name = "".join(c if c.isalnum() else "_" for c in config.name)
        pptx_filename = f"smart_discovery_{safe_name}.pptx"
        pptx_path = OUTPUT_DIR / pptx_filename
        prs.save(str(pptx_path))

        return pptx_path, slides_created

    except Exception as e:
        logger.error(f"PowerPoint generation failed: {e}")
        return None, 0


# ============================================================================
# Legacy Entry Point (backward compatibility)
# ============================================================================
def smart_discover_legacy(
    notion_url: str = None,
    raw_content: str = None,
    customer_name: str = "Client",
) -> Dict[str, Any]:
    """Legacy entry point without CustomerConfig."""
    config = CustomerConfig(name=customer_name)
    return smart_discover(notion_url, raw_content, customer_name, config)


# ============================================================================
# Main
# ============================================================================
if __name__ == "__main__":
    # Test with sample content
    test_content = """
    # Product Interview - TaskFlow App

    Interviewed Sarah Chen, Product Owner, on January 15, 2025.

    ## Product Vision
    TaskFlow aims to be the simplest task management app for small teams of 2-10 people.
    We want to help teams stay organized without the complexity of enterprise tools.

    ## Target Users
    - Freelancers managing multiple clients
    - Small agency teams
    - Remote workers needing simple collaboration

    ## Pain Points
    - "Existing tools are too complex"
    - "Trello boards get messy with a team"
    - "We use spreadsheets because nothing else is simple enough"

    ## Key Features
    1. One-click task creation
    2. Team assignment with @mentions
    3. Due date tracking

    ## Business Model
    - Free tier for individuals
    - $5/user/month for teams
    - Target: 10,000 paying teams in Year 1

    ## Success Metrics
    - User activation rate: > 60%
    - Weekly retention: > 40%
    - NPS score: > 50
    """

    print("Testing Smart Discovery (Production)...")

    config = CustomerConfig(
        name="TaskFlow",
        must_include=["Executive Summary", "Key Findings", "Recommendations"],
        slide_budget={"min": 8, "max": 30, "per_section_max": 4}
    )

    result = smart_discover(
        raw_content=test_content,
        customer_name="TaskFlow",
        config=config
    )

    print("\n" + "=" * 60)
    print("RESULTS")
    print("=" * 60)
    print(f"Success: {result['success']}")
    print(f"Evidence items: {result['evidence_count']}")
    print(f"Sections: {len(result['sections'])}")
    for s in result['sections']:
        print(f"  - {s}")
    print(f"Validation: {result['validation']}")
    print(f"Markdown: {result['markdown_path']}")
    print(f"PowerPoint: {result['powerpoint_path']}")
    if result['error']:
        print(f"Error: {result['error']}")
