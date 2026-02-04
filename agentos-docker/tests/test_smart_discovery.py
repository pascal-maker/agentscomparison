"""
Tests for Smart Discovery (Production Grade)
============================================

Test evidence grounding, slide generation, and validation.
"""

import os
import sys
import re
from pathlib import Path

# Add parent to path
sys.path.insert(0, str(Path(__file__).parent.parent))

from dotenv import load_dotenv
load_dotenv()

import pytest

from shared.evidence import (
    EvidenceItem,
    EvidenceCollection,
    GroundedSection,
    GroundedBullet,
    CustomerConfig,
    extract_evidence_from_content,
    validate_grounded_report
)
from teams.smart_discovery import smart_discover


# ============================================================================
# Test Data
# ============================================================================

TECHNICAL_INPUT = """
# Technical Discovery: Countroll Platform

Interview with Alexander, CTO - December 18, 2024

## Platform Overview
Countroll is a cloud-based roll asset management platform for tracking industrial rollers.
The system monitors roller lifecycle from installation to replacement.

## Technical Architecture
- Web portal: Vue.js frontend
- Mobile app: Android using Kotlin Multiplatform with MVI pattern
- Backend: Kotlin with Spring Boot framework
- Database: Azure Cosmos DB for domain data

## Infrastructure
- Hosted on Microsoft Azure
- Uses Application Gateway for routing
- Load balancing across multiple VMs

## Development Practices
- Agile process using Shortcut for backlog management
- CI/CD with Azure DevOps and GitHub Actions
- Code coverage target: 50-70%
- QA managed through Qase platform

## Support
- Jira Service Desk for external customer support
- Internal IT handled separately
- SLA tracking for customer issues
"""

PRODUCT_INPUT = """
# Product Discovery: TaskFlow App

Interview with Sarah Chen, Product Owner - January 15, 2025

## Product Vision
TaskFlow aims to be the simplest task management app for small teams of 2-10 people.

## Target Users
- Freelancers managing multiple clients
- Small agency teams (design, marketing)
- Remote workers needing collaboration

## Pain Points from Research
- "Existing tools are too complex - I spend more time managing the tool"
- "Trello boards get messy fast"
- "Notion requires too much setup"

## Competitive Landscape
- Todoist: Good for individuals, lacks team features
- Trello: Boards don't work for everyone
- Linear: Too developer-focused

## Key Features
1. One-click task creation
2. Team assignment with @mentions
3. Due date tracking with calendar
4. Mobile app (iOS priority)

## Business Model
- Free tier: Individuals, up to 3 projects
- Team tier: $5/user/month, unlimited
- Target: 10,000 paying teams Year 1

## Success Metrics
- User activation: > 60%
- Weekly retention: > 40%
- NPS: > 50
"""


# ============================================================================
# Test: Evidence Extraction
# ============================================================================

def test_evidence_extraction():
    """Test that evidence is extracted with IDs."""
    evidence = extract_evidence_from_content(
        TECHNICAL_INPUT,
        page_title="Countroll Discovery",
        page_id="test-123"
    )

    assert len(evidence) > 0, "Should extract evidence items"

    # Check evidence items have IDs
    for item in evidence.items.values():
        assert item.id.startswith("EVID-"), f"Evidence ID should start with EVID-: {item.id}"
        assert len(item.quote) > 0, "Evidence should have quote"


def test_evidence_search():
    """Test evidence search functionality."""
    evidence = extract_evidence_from_content(TECHNICAL_INPUT)

    # Search for Azure-related evidence
    results = evidence.search(["Azure", "cloud"])
    assert len(results) > 0, "Should find Azure-related evidence"


# ============================================================================
# Test: Grounded Sections
# ============================================================================

def test_grounded_bullet_validation():
    """Test that bullets without evidence are flagged."""
    section = GroundedSection(name="Test Section")

    # Add grounded bullet
    section.add_bullet("Grounded point", ["EVID-abc123"])

    # Add ungrounded bullet
    section.bullets.append(GroundedBullet(text="Ungrounded point", evidence_ids=[]))

    # Validate
    ungrounded = section.validate_grounding()
    assert len(ungrounded) == 1, "Should detect 1 ungrounded bullet"
    assert "Ungrounded point" in ungrounded[0]


def test_section_with_only_open_questions():
    """Test section with no evidence creates open questions."""
    section = GroundedSection(name="Empty Section")

    # No bullets, just open questions
    section.add_open_question("What is the timeline?")
    section.add_open_question("What is the budget?")

    assert len(section.open_questions) == 2
    assert not section.has_sufficient_evidence


# ============================================================================
# Test: Customer Config
# ============================================================================

def test_customer_config_must_include():
    """Test that missing must_include concepts are detected."""
    config = CustomerConfig(
        name="Test",
        must_include=["Executive Summary", "Recommendations", "Budget Analysis"]
    )

    sections = ["Executive Summary", "Technical Overview"]
    missing = config.validate_sections(sections)

    assert "Recommendations" in missing
    assert "Budget Analysis" in missing
    assert "Executive Summary" not in missing


def test_customer_config_terminology():
    """Test terminology mapping."""
    config = CustomerConfig(
        name="Test",
        terminology_map={"client": "Acme Corp", "the customer": "Acme Corp"}
    )

    text = "The client needs a solution. The customer will benefit."
    result = config.apply_terminology(text)

    assert "Acme Corp" in result
    assert "client" not in result.lower() or "Acme Corp" in result


def test_slide_budget():
    """Test slide budget constraints."""
    config = CustomerConfig(
        name="Test",
        slide_budget={"min": 10, "max": 30, "per_section_max": 4}
    )

    assert config.get_min_slides() == 10
    assert config.get_max_slides() == 30
    assert config.get_per_section_max() == 4


# ============================================================================
# Test: Smart Discovery - Technical Input
# ============================================================================

def test_smart_discovery_technical():
    """Test smart discovery with technical input."""
    config = CustomerConfig(
        name="Countroll",
        must_include=["Executive Summary", "Key Findings", "Recommendations"],
        slide_budget={"min": 8, "max": 40, "per_section_max": 6}
    )

    result = smart_discover(
        raw_content=TECHNICAL_INPUT,
        customer_name="Countroll",
        config=config
    )

    # Basic success
    assert result["success"], f"Should succeed: {result.get('error')}"
    assert result["evidence_count"] > 0, "Should have evidence"
    assert len(result["sections"]) > 0, "Should have sections"

    # Check files exist
    assert result["markdown_path"], "Should have markdown path"
    assert result["powerpoint_path"], "Should have powerpoint path"
    assert Path(result["markdown_path"]).exists(), "Markdown file should exist"
    assert Path(result["powerpoint_path"]).exists(), "PowerPoint file should exist"

    # Check markdown has evidence references
    md_content = Path(result["markdown_path"]).read_text()
    assert "[EVID-" in md_content, "Markdown should have evidence references"

    # Check has agenda
    assert "## Agenda" in md_content, "Should have Agenda section"

    # Check has open questions somewhere (at least one section should have gaps)
    # Note: This may not always be true if all info is present
    # assert "Open Questions" in md_content or result["validation"]["valid"]


def test_smart_discovery_product():
    """Test smart discovery with product input."""
    config = CustomerConfig(
        name="TaskFlow",
        must_include=["Executive Summary", "Key Findings", "Recommendations"],
        slide_budget={"min": 8, "max": 40, "per_section_max": 6}
    )

    result = smart_discover(
        raw_content=PRODUCT_INPUT,
        customer_name="TaskFlow",
        config=config
    )

    # Basic success
    assert result["success"], f"Should succeed: {result.get('error')}"
    assert result["evidence_count"] > 0, "Should have evidence"

    # Check evidence grounding
    md_content = Path(result["markdown_path"]).read_text()

    # Count evidence references
    evid_count = len(re.findall(r'\[EVID-[a-f0-9]+\]', md_content))
    assert evid_count > 0, "Should have evidence references in markdown"


# ============================================================================
# Test: PowerPoint Structure
# ============================================================================

def test_powerpoint_has_title_and_agenda():
    """Test that PowerPoint has required slides."""
    from pptx import Presentation

    config = CustomerConfig(
        name="TestPPT",
        slide_budget={"min": 5, "max": 20, "per_section_max": 3}
    )

    result = smart_discover(
        raw_content=TECHNICAL_INPUT,
        customer_name="TestPPT",
        config=config
    )

    assert result["powerpoint_path"], "Should have PowerPoint"

    # Open and check structure
    prs = Presentation(result["powerpoint_path"])

    assert len(prs.slides) >= 2, "Should have at least Title + Agenda slides"

    # Check first slide has title content
    first_slide = prs.slides[0]
    has_title_text = False
    for shape in first_slide.shapes:
        if shape.has_text_frame:
            if "TestPPT" in shape.text_frame.text or "Discovery" in shape.text_frame.text:
                has_title_text = True
                break
    assert has_title_text, "First slide should have title"

    # Check second slide is Agenda
    second_slide = prs.slides[1]
    has_agenda = False
    for shape in second_slide.shapes:
        if shape.has_text_frame:
            if "Agenda" in shape.text_frame.text:
                has_agenda = True
                break
    assert has_agenda, "Second slide should be Agenda"


def test_powerpoint_speaker_notes():
    """Test that PowerPoint has evidence in speaker notes."""
    from pptx import Presentation

    config = CustomerConfig(name="NotesTest")

    result = smart_discover(
        raw_content=TECHNICAL_INPUT,
        customer_name="NotesTest",
        config=config
    )

    prs = Presentation(result["powerpoint_path"])

    # Check that at least some slides have speaker notes
    slides_with_notes = 0
    for slide in prs.slides:
        notes = slide.notes_slide
        if notes and notes.notes_text_frame.text:
            slides_with_notes += 1

    assert slides_with_notes > 0, "Should have slides with speaker notes"


# ============================================================================
# Test: Validation
# ============================================================================

def test_grounding_validation():
    """Test that validation catches ungrounded bullets."""
    sections = [
        GroundedSection(name="Good Section"),
        GroundedSection(name="Bad Section"),
    ]

    # Good section has grounded bullets
    sections[0].add_bullet("Point with evidence", ["EVID-abc123"])

    # Bad section has ungrounded bullet
    sections[1].bullets.append(GroundedBullet(text="No evidence", evidence_ids=[]))

    validation = validate_grounded_report(sections)

    assert not validation["valid"], "Should be invalid with ungrounded bullets"
    assert validation["ungrounded_bullets"], "Should list ungrounded bullets"


# ============================================================================
# Test: Slide Budget Enforcement
# ============================================================================

def test_slide_budget_max_enforcement():
    """Test that slide budget maximum is enforced."""
    from agents.revisor import enforce_slide_budget

    # Create many sections with many bullets
    sections = []
    for i in range(10):
        section = GroundedSection(name=f"Section {i}")
        for j in range(20):  # 20 bullets per section
            section.add_bullet(f"Point {j}", [f"EVID-{i}{j}"])
        sections.append(section)

    config = CustomerConfig(
        name="Test",
        slide_budget={"min": 5, "max": 20, "per_section_max": 2}
    )

    # This would create way more than 20 slides without enforcement
    enforced = enforce_slide_budget(sections, config)

    # Check bullets were trimmed
    total_bullets = sum(len(s.bullets) for s in enforced)
    assert total_bullets < 200, "Should have trimmed bullets to fit budget"


# ============================================================================
# Main
# ============================================================================

if __name__ == "__main__":
    print("Running Smart Discovery Tests...")
    print("=" * 60)

    # Run tests manually for visibility
    tests = [
        test_evidence_extraction,
        test_evidence_search,
        test_grounded_bullet_validation,
        test_section_with_only_open_questions,
        test_customer_config_must_include,
        test_customer_config_terminology,
        test_slide_budget,
        test_smart_discovery_technical,
        test_smart_discovery_product,
        test_powerpoint_has_title_and_agenda,
        test_powerpoint_speaker_notes,
        test_grounding_validation,
        test_slide_budget_max_enforcement,
    ]

    passed = 0
    failed = 0

    for test in tests:
        try:
            test()
            print(f"  ✓ {test.__name__}")
            passed += 1
        except AssertionError as e:
            print(f"  ✗ {test.__name__}: {e}")
            failed += 1
        except Exception as e:
            print(f"  ✗ {test.__name__}: {type(e).__name__}: {e}")
            failed += 1

    print("=" * 60)
    print(f"Results: {passed} passed, {failed} failed")
